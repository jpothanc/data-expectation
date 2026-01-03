#!/usr/bin/env python3
"""
Script to convert PRESENTATION.md to PowerPoint format.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re

def parse_markdown_slides(md_file):
    """Parse markdown file and extract slides."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by slide separators (---)
    slides = []
    current_slide = []
    
    for line in content.split('\n'):
        if line.strip() == '---':
            if current_slide:
                slides.append('\n'.join(current_slide))
                current_slide = []
        else:
            current_slide.append(line)
    
    if current_slide:
        slides.append('\n'.join(current_slide))
    
    return slides

def parse_slide_content(slide_text):
    """Extract title and content from slide text."""
    lines = slide_text.strip().split('\n')
    title = None
    content = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Check for title (starts with #)
        if line.startswith('#') and not title:
            title = line.lstrip('#').strip()
        elif line.startswith('##'):
            if not title:
                title = line.lstrip('#').strip()
            else:
                content.append(('subtitle', line.lstrip('#').strip()))
        elif line.startswith('-') or line.startswith('*'):
            content.append(('bullet', line.lstrip('-*').strip()))
        elif line.startswith('```'):
            continue  # Skip code blocks for now
        elif line.startswith('`'):
            content.append(('code', line.strip('`')))
        else:
            content.append(('text', line))
    
    return title or 'Slide', content

def create_presentation(md_file, output_file):
    """Create PowerPoint from markdown."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slides = parse_markdown_slides(md_file)
    
    for slide_text in slides:
        if not slide_text.strip():
            continue
        
        title, content = parse_slide_content(slide_text)
        
        # Choose layout
        if 'Title Slide' in title or slides.index(slide_text) == 0:
            slide_layout = prs.slide_layouts[0]  # Title slide
        else:
            slide_layout = prs.slide_layouts[1]  # Title and content
        
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Add content
        if hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.word_wrap = True
            
            for content_type, text in content:
                if content_type == 'bullet':
                    p = tf.add_paragraph()
                    p.text = text
                    p.level = 0
                    p.font.size = Pt(14)
                elif content_type == 'subtitle':
                    p = tf.add_paragraph()
                    p.text = text
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.bold = True
                elif content_type == 'text':
                    p = tf.add_paragraph()
                    p.text = text
                    p.level = 0
                    p.font.size = Pt(12)
                elif content_type == 'code':
                    p = tf.add_paragraph()
                    p.text = text
                    p.level = 0
                    p.font.size = Pt(10)
                    p.font.name = 'Courier New'
    
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == '__main__':
    import sys
    
    md_file = 'PRESENTATION.md'
    output_file = 'Data_Expectations_Presentation.pptx'
    
    if len(sys.argv) > 1:
        md_file = sys.argv[1]
    if len(sys.argv) > 2:
        output_file = sys.argv[2]
    
    try:
        create_presentation(md_file, output_file)
        print(f"\n✅ Successfully created PowerPoint: {output_file}")
        print(f"📊 Total slides: {len(parse_markdown_slides(md_file))}")
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("   Install with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()

